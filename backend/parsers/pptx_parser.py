from pptx import Presentation


def parse_pptx(file_path):
    """
    Extract text from PowerPoint slides.
    """

    presentation = Presentation(file_path)

    extracted_text = []

    for slide_number, slide in enumerate(presentation.slides, start=1):

        extracted_text.append(f"\nSlide Number: {slide_number}\n")

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                extracted_text.append(shape.text)

    return "\n".join(extracted_text)